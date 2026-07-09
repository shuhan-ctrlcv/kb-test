#!/usr/bin/env python3
"""Generate binary gap-probe fixtures for the PedalWorks corpus.

Creates four binary files (PDF, DOCX, XLSX, PPTX) that exercise the real
extractors in bitwise_slack_bot's app/ingest/preprocess/extractors.py. All
content is PedalWorks-specific -- it cross-references the same suppliers,
parts, and systems described in docs/ and eval/expected_structure.yaml -- so
downstream eval questions can target it.

The PDF embeds one fact ONLY inside a PIL-drawn image (never in the page's
extractable text), and the DOCX puts QC standards ONLY inside a table (never
in a paragraph). Both are known extractor gaps: images are dropped when
extracting PDF text, and only paragraph text -- not table cells -- is read
from a DOCX. This script documents those gaps as fixtures; it does not fix
them.

Run from the kb-test repo root:
    /home/nine/anaconda3/envs/road_asset/bin/python tools/gen_binary_fixtures.py
"""
from __future__ import annotations

import io
from pathlib import Path

import fitz  # PyMuPDF
import openpyxl
from docx import Document
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent.parent

# Supplier facts, consistent with docs/procurement.md and docs/trailblazer.md.
PARTS = [
    # part, supplier, lead_time_days, unit_cost
    ("Frame", "FrameForge Ltd", 14, 120),
    ("Wheelset", "RollRight Co", 7, 85),
    ("Drivetrain", "GearWorks Inc", 10, 95),
    ("Brake Set", "BrakeSafe GmbH", 10, 40),
    ("Suspension Fork", "ForkFactory Co", 12, 150),
]

# The one fact that exists ONLY inside the PDF's embedded image, never in its
# extractable text -- the gap probe for "images are dropped on PDF extract".
BAY_FACT = "Suspension Fork inventory is staged in Bay 12 of the North Intake Warehouse."


def make_bay_map_image() -> bytes:
    """A small PIL-drawn 'warehouse bay map' diagram. BAY_FACT is captioned
    on the image itself and must not appear anywhere in the PDF's text."""
    img = Image.new("RGB", (620, 300), "white")
    draw = ImageDraw.Draw(img)
    font = ImageFont.load_default()
    draw.rectangle([10, 10, 610, 290], outline="black", width=2)
    draw.text((20, 20), "North Intake Warehouse -- Bay Map", fill="black", font=font)
    for i, bay in enumerate(range(1, 13)):
        col, row = i % 6, i // 6
        x, y = 20 + col * 98, 55 + row * 95
        outline = "red" if bay == 12 else "gray"
        width = 3 if bay == 12 else 1
        draw.rectangle([x, y, x + 85, y + 80], outline=outline, width=width)
        draw.text((x + 5, y + 5), f"Bay {bay}", fill="black", font=font)
    draw.text((20, 260), BAY_FACT, fill="black", font=font)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


def build_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page()
    y = 50

    def line(text: str, size: float = 11, dy: float = 18) -> None:
        nonlocal y
        page.insert_text((50, y), text, fontsize=size)
        y += dy

    line("PedalWorks Supplier Catalog", size=16, dy=28)
    line("Prices and lead times for the five purchased parts behind the", size=11)
    line("City Cruiser and Trail Blazer, received at the North Intake Warehouse.", size=11, dy=26)

    line("Part             Supplier              Lead Time    Unit Cost", size=11, dy=18)
    for part, supplier, lead, cost in PARTS:
        line(f"{part:<17}{supplier:<22}{lead} days       ${cost}", size=11)
    y += 20
    line("See the bay map below for warehouse storage assignment.", size=11, dy=24)

    img_bytes = make_bay_map_image()
    rect = fitz.Rect(50, y, 50 + 465, y + 225)
    page.insert_image(rect, stream=img_bytes)

    doc.save(str(path))
    doc.close()


def build_docx(path: Path) -> None:
    doc = Document()
    doc.add_heading("PedalWorks Quality Checklist", level=1)
    doc.add_paragraph(
        "This checklist enumerates the gate criteria used at the Portland "
        "Assembly Floor's QC station before a City Cruiser or Trail Blazer "
        "is released to Shipping. Every inspection -- pass or fail -- "
        "consumes test material and posts a scrap cost to the Ledger."
    )
    doc.add_paragraph(
        "The first unit inspected in a run always undergoes a first-article "
        "inspection and always fails; it is returned to Assembly for rework "
        "and re-inspected before moving on."
    )

    table = doc.add_table(rows=1, cols=4)
    hdr = table.rows[0].cells
    hdr[0].text, hdr[1].text, hdr[2].text, hdr[3].text = (
        "Check Item",
        "Applies To",
        "Standard",
        "Pass Threshold",
    )
    rows = [
        ("Frame weld integrity", "City Cruiser, Trail Blazer", "ISO 4210 weld spec", "No visible cracks"),
        ("Wheelset trueness", "City Cruiser, Trail Blazer", "Lateral runout", "<= 0.5 mm"),
        ("Drivetrain shift index", "City Cruiser", "10-speed indexed shift", "All gears engage cleanly"),
        ("Brake Set pull test", "City Cruiser", "Lever force", "Holds at 800 N, no pad slip"),
        ("Suspension Fork travel test", "Trail Blazer", "Compression/rebound", "Full 120 mm travel"),
    ]
    for item, applies, standard, threshold in rows:
        cells = table.add_row().cells
        cells[0].text, cells[1].text, cells[2].text, cells[3].text = item, applies, standard, threshold

    doc.add_paragraph(
        "Any unit failing one or more checklist items above is routed back "
        "to Assembly for rework, consistent with PedalWorks' QC rework loop."
    )
    doc.save(str(path))


def build_xlsx(path: Path) -> None:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Stock"
    ws.append(["Part", "Supplier", "Warehouse", "Qty On Hand", "Reorder Threshold", "Unit Cost"])
    stock = [
        ("Frame", "FrameForge Ltd", "North Intake Warehouse", 340, 100, 120),
        ("Wheelset", "RollRight Co", "North Intake Warehouse", 512, 150, 85),
        ("Drivetrain", "GearWorks Inc", "North Intake Warehouse", 208, 80, 95),
        ("Brake Set", "BrakeSafe GmbH", "North Intake Warehouse", 275, 90, 40),
        ("Suspension Fork", "ForkFactory Co", "North Intake Warehouse", 64, 50, 150),
    ]
    for row in stock:
        ws.append(row)
    wb.save(str(path))


def build_pptx(path: Path) -> None:
    pres = Presentation()

    slide = pres.slides.add_slide(pres.slide_layouts[0])  # Title Slide
    slide.shapes.title.text = "PedalWorks Factory Overview"
    slide.placeholders[1].text = "From supplier to shipped bicycle"

    slide = pres.slides.add_slide(pres.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Two Hubs"
    body = slide.placeholders[1].text_frame
    body.text = "Scheduler -- fan-out hub: reads PartsDB + OrderDB, triggers Procurement, releases Assembly work"
    body.add_paragraph().text = "Ledger -- fan-in hub: every stage posts a cost here, backed by LedgerDB"

    slide = pres.slides.add_slide(pres.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Three Systems of Record"
    body = slide.placeholders[1].text_frame
    body.text = "PartsDB -- current stock level of each part"
    body.add_paragraph().text = "OrderDB -- customer orders, open or shipped"
    body.add_paragraph().text = "LedgerDB -- storage backing the Ledger"

    slide = pres.slides.add_slide(pres.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "Supplier Lead Times"
    rows, cols = len(PARTS) + 1, 3
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(3)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text, table.cell(0, 1).text, table.cell(0, 2).text = "Part", "Supplier", "Lead Time"
    for r, (part, supplier, lead, _cost) in enumerate(PARTS, start=1):
        table.cell(r, 0).text = part
        table.cell(r, 1).text = supplier
        table.cell(r, 2).text = f"{lead} days"

    pres.save(str(path))


def main() -> None:
    (ROOT / "docs").mkdir(parents=True, exist_ok=True)
    (ROOT / "data").mkdir(parents=True, exist_ok=True)
    build_pdf(ROOT / "docs" / "supplier_catalog.pdf")
    build_docx(ROOT / "docs" / "quality_checklist.docx")
    build_xlsx(ROOT / "data" / "inventory_stock.xlsx")
    build_pptx(ROOT / "docs" / "factory_overview.pptx")
    print("Generated binary fixtures under docs/ and data/.")


if __name__ == "__main__":
    main()
